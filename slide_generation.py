"""
Functions to create and manipulate PowerPoint presentations using python-pptx.
Extended to implement strict slide template with placeholders for logos and visuals.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_presentation():
    prs = Presentation()
    return prs

def add_cover_slide(prs, client_name, project_name="Proposal"):
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = client_name
    slide.placeholders[1].text = f"Proposal for {project_name}"
    # Placeholder for logo can be added here

def add_title_slide(prs, title, subtitle):
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle

def add_content_slide(prs, title, bullet_points):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    text_frame = slide.placeholders[1].text_frame
    text_frame.clear()
    for point in bullet_points:
        p = text_frame.add_paragraph()
        p.text = point
        p.level = 0

def add_horizontal_roadmap_slide(prs, title, phases):
    """
    Add a slide with a horizontal roadmap layout showing phases with titles and descriptions.
    """
    slide_layout = prs.slide_layouts[1]  # Title and Content layout with placeholder idx 1
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title

    shapes = slide.shapes
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(9)
    height = Inches(4.5)

    # Add a simple horizontal list of phases as bullet points for now
    text_frame = shapes.placeholders[1].text_frame
    text_frame.clear()
    for phase in phases:
        title = phase.get("title", "")
        desc = phase.get("description", "")
        p = text_frame.add_paragraph()
        p.text = f"{title}: {desc}"
        p.level = 0

def save_presentation(prs, file_path):
    prs.save(file_path)
